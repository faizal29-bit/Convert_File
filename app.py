"""
LocalConvert - Sistem konversi file lokal
Semua proses berjalan di komputer sendiri, tanpa limit
"""

from flask import Flask, request, jsonify, send_file, render_template
import os, uuid, subprocess, shutil, platform, tempfile, threading, time, zipfile
from pathlib import Path
from werkzeug.utils import secure_filename

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = None  # Tanpa limit untuk lokal

TEMP_DIR = Path(tempfile.gettempdir()) / 'localconvert_jobs'
TEMP_DIR.mkdir(exist_ok=True)

jobs = {}  # job_id -> {dir, files, created}

# ─── HELPER: LibreOffice ───────────────────────────────────────────────────

def get_libreoffice():
    system = platform.system()
    if system == 'Windows':
        candidates = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            r'C:\Program Files\LibreOffice 7\program\soffice.exe',
            r'C:\Program Files\LibreOffice 24\program\soffice.exe',
        ]
        for c in candidates:
            if os.path.exists(c):
                return c
        return shutil.which('soffice')
    elif system == 'Darwin':
        candidates = ['/Applications/LibreOffice.app/Contents/MacOS/soffice']
        for c in candidates:
            if os.path.exists(c):
                return c
        return shutil.which('libreoffice') or shutil.which('soffice')
    else:
        return shutil.which('libreoffice') or shutil.which('soffice')

def run_libreoffice(input_path, output_format, output_dir):
    lo = get_libreoffice()
    if not lo:
        raise RuntimeError(
            "LibreOffice tidak ditemukan.\n"
            "Install dari: https://www.libreoffice.org/download/libreoffice/"
        )
    # LibreOffice needs a user profile dir to run headless without conflicts
    user_dir = TEMP_DIR / 'lo_user'
    user_dir.mkdir(exist_ok=True)

    cmd = [
        lo, '--headless',
        f'-env:UserInstallation=file:///{user_dir.as_posix()}',
        '--convert-to', output_format,
        '--outdir', str(output_dir),
        str(input_path)
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice gagal: {result.stderr or result.stdout}")

    stem = Path(input_path).stem
    expected = output_dir / f"{stem}.{output_format}"
    if not expected.exists():
        # Try to find any matching output
        matches = list(output_dir.glob(f"{stem}*"))
        if matches:
            return matches[0]
        raise RuntimeError("File output tidak ditemukan setelah konversi LibreOffice")
    return expected

def get_poppler_path():
    """Cari poppler untuk pdf2image di Windows"""
    if platform.system() == 'Windows':
        candidates = [
            r'C:\poppler\bin',
            r'C:\Program Files\poppler\bin',
            r'C:\poppler-windows\bin',
        ]
        for c in candidates:
            if os.path.exists(c):
                return c
        # Try PATH
        pdftoppm = shutil.which('pdftoppm')
        if pdftoppm:
            return str(Path(pdftoppm).parent)
    return None

# ─── KONVERSI: Gambar → PDF ───────────────────────────────────────────────

def convert_images_to_pdf(input_files, output_path):
    import img2pdf
    from PIL import Image
    processed, temps = [], []
    try:
        for f in input_files:
            img = Image.open(str(f))
            # Cek jika gambar memiliki saluran transparan (Alpha)
            if img.mode in ('RGBA', 'LA', 'PA') or 'transparency' in img.info:
                tmp = str(f) + '_cvt.png'
                img.save(tmp, 'PNG')
            else:
                if img.mode not in ('RGB', 'L'):
                    img = img.convert('RGB')
                tmp = str(f) + '_cvt.jpg'
                img.save(tmp, 'JPEG', quality=95)
            
            processed.append(tmp)
            temps.append(tmp)
            
        with open(str(output_path), 'wb') as out:
            out.write(img2pdf.convert(processed))
    finally:
        for t in temps:
            try: os.remove(t)
            except: pass

# ─── KONVERSI: Gambar → Word ──────────────────────────────────────────────

def convert_images_to_word(input_files, output_path):
    from docx import Document
    from docx.shared import Inches, Pt
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from PIL import Image

    doc = Document()
    # Hapus margin bawaan yang terlalu besar
    section = doc.sections[0]
    section.top_margin = Inches(0.5)
    section.bottom_margin = Inches(0.5)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)

    for i, img_path in enumerate(input_files):
        img = Image.open(str(img_path))
        w, h = img.size

        # Hitung lebar agar muat di halaman (max 6.5 inci)
        max_w_in = 6.5
        w_in = w / 96  # asumsi 96 DPI
        h_in = h / 96

        if w_in > max_w_in:
            scale = max_w_in / w_in
            w_in = max_w_in
            h_in = h_in * scale

        # Jika tinggi melebihi halaman, scale lagi
        max_h_in = 9.0
        if h_in > max_h_in:
            scale = max_h_in / h_in
            w_in = w_in * scale
            h_in = max_h_in

        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = para.add_run()
        run.add_picture(str(img_path), width=Inches(w_in))

        if i < len(input_files) - 1:
            doc.add_page_break()

    doc.save(str(output_path))

# ─── KONVERSI: Word → PDF ────────────────────────────────────────────────

def convert_word_to_pdf(input_file, output_path):
    output_dir = Path(output_path).parent
    result = run_libreoffice(input_file, 'pdf', output_dir)
    expected = output_dir / (Path(input_file).stem + '.pdf')
    if result != Path(output_path) and expected != Path(output_path):
        shutil.move(str(result), str(output_path))
    elif result != Path(output_path):
        shutil.move(str(result), str(output_path))

# ─── KONVERSI: PDF → Word ────────────────────────────────────────────────

def convert_pdf_to_word(input_file, output_path):

    import os
    import io
    import re
    import subprocess
    import fitz

    from collections import Counter

    from pdf2docx import Converter

    from docx import Document

    from docx.shared import (
        Pt,
        RGBColor,
        Cm,
        Inches
    )

    from docx.enum.text import (
        WD_ALIGN_PARAGRAPH
    )

    from docx.oxml.ns import qn

    from docx.oxml import OxmlElement

    # =====================================================
    # VALIDASI
    # =====================================================

    if not os.path.exists(input_file):

        print("PDF tidak ditemukan!")
        return

    # =====================================================
    # METHOD 1
    # PDF2DOCX ENGINE
    # =====================================================

    try:

        print("=" * 60)
        print("METHOD 1 : PDF2DOCX")
        print("=" * 60)

        cv = Converter(str(input_file))

        cv.convert(
            str(output_path),
            start=0,
            end=None,
            multi_processing=True
        )

        cv.close()

        if os.path.exists(output_path):

            print("SUCCESS PDF2DOCX")
            return

    except Exception as e:

        print("PDF2DOCX gagal:")
        print(e)

    # =====================================================
    # METHOD 2
    # LIBREOFFICE
    # =====================================================

    try:

        print("=" * 60)
        print("METHOD 2 : LIBREOFFICE")
        print("=" * 60)

        output_dir = os.path.dirname(
            os.path.abspath(output_path)
        )

        cmd = [

            "soffice",

            "--headless",

            "--convert-to",
            "docx",

            str(input_file),

            "--outdir",
            output_dir
        ]

        subprocess.run(
            cmd,
            check=True
        )

        if os.path.exists(output_path):

            print("SUCCESS LIBREOFFICE")
            return

    except Exception as e:

        print("LIBREOFFICE gagal:")
        print(e)

    # =====================================================
    # METHOD 3
    # ADVANCED RECONSTRUCTION
    # =====================================================

    try:

        print("=" * 60)
        print("METHOD 3 : ADVANCED RECONSTRUCTION")
        print("=" * 60)

        # =================================================
        # DOCUMENT
        # =================================================

        doc = Document()

        # =================================================
        # GLOBAL STYLE
        # =================================================

        style = doc.styles['Normal']

        style.font.name = 'Times New Roman'

        style._element.rPr.rFonts.set(
            qn('w:eastAsia'),
            'Times New Roman'
        )

        style.font.size = Pt(12)

        # =================================================
        # MARGIN
        # =================================================

        for section in doc.sections:

            section.top_margin = Cm(2.54)
            section.bottom_margin = Cm(2.54)

            section.left_margin = Cm(2.54)
            section.right_margin = Cm(2.54)

        # =================================================
        # OPEN PDF
        # =================================================

        pdf = fitz.open(str(input_file))

        # =================================================
        # HELPER
        # =================================================

        def clean_text(text):

            if not text:
                return ""

            text = text.replace("\n", " ")

            text = re.sub(
                r'\s+',
                ' ',
                text
            )

            text = re.sub(
                r'\s+([.,!?;:])',
                r'\1',
                text
            )

            return text.strip()

        def get_rgb(color_int):

            r = (color_int >> 16) & 255
            g = (color_int >> 8) & 255
            b = color_int & 255

            return r, g, b

        def set_cell_border(cell):

            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()

            tcBorders = OxmlElement('w:tcBorders')

            for edge in (
                'top',
                'left',
                'bottom',
                'right'
            ):

                edge_element = OxmlElement(f'w:{edge}')

                edge_element.set(
                    qn('w:val'),
                    'single'
                )

                edge_element.set(
                    qn('w:sz'),
                    '4'
                )

                edge_element.set(
                    qn('w:space'),
                    '0'
                )

                edge_element.set(
                    qn('w:color'),
                    '000000'
                )

                tcBorders.append(edge_element)

            tcPr.append(tcBorders)

        # =================================================
        # DETECT BODY FONT
        # =================================================

        all_sizes = []

        for page in pdf:

            blocks = page.get_text(
                "dict",
                flags=(
                    fitz.TEXT_PRESERVE_WHITESPACE
                    | fitz.TEXT_PRESERVE_LIGATURES
                )
            )["blocks"]

            for block in blocks:

                if block.get("type") != 0:
                    continue

                for line in block.get("lines", []):

                    for span in line.get("spans", []):

                        size = span.get(
                            "size",
                            12
                        )

                        if size > 5:

                            all_sizes.append(
                                round(size)
                            )

        if all_sizes:

            common_size = Counter(
                all_sizes
            ).most_common(1)[0][0]

        else:

            common_size = 12

        # =================================================
        # PAGE LOOP
        # =================================================

        for page_index in range(len(pdf)):

            page = pdf[page_index]

            if page_index > 0:
                doc.add_page_break()

            page_width = page.rect.width

            # =============================================
            # TABLE DETECTION
            # =============================================

            try:

                tables = page.find_tables()

            except:

                tables = []

            # =============================================
            # BLOCKS
            # =============================================

            blocks = page.get_text(
                "dict",
                flags=(
                    fitz.TEXT_PRESERVE_WHITESPACE
                    | fitz.TEXT_PRESERVE_LIGATURES
                )
            )["blocks"]

            blocks.sort(
                key=lambda b: (
                    b["bbox"][1],
                    b["bbox"][0]
                )
            )

            previous_y = None

            # =============================================
            # BLOCK LOOP
            # =============================================

            for block in blocks:

                # =========================================
                # IMAGE
                # =========================================

                if block["type"] == 1:

                    try:

                        img_data = pdf.extract_image(
                            block["xref"]
                        )

                        image_bytes = img_data["image"]

                        doc.add_picture(
                            io.BytesIO(image_bytes),
                            width=Inches(5.8)
                        )

                        last_paragraph = (
                            doc.paragraphs[-1]
                        )

                        last_paragraph.alignment = (
                            WD_ALIGN_PARAGRAPH.CENTER
                        )

                        last_paragraph.paragraph_format.space_after = Pt(10)

                    except:
                        pass

                    continue

                # =========================================
                # SKIP NON TEXT
                # =========================================

                if block["type"] != 0:
                    continue

                lines = block.get(
                    "lines",
                    []
                )

                if not lines:
                    continue

                # =========================================
                # BLOCK INFO
                # =========================================

                bx0, by0, bx1, by1 = (
                    block["bbox"]
                )

                block_width = bx1 - bx0

                # =========================================
                # PARAGRAPH TEXT
                # =========================================

                paragraph_lines = []

                for line in lines:

                    line_text = ""

                    for span in line.get(
                        "spans",
                        []
                    ):

                        line_text += span.get(
                            "text",
                            ""
                        )

                    line_text = clean_text(
                        line_text
                    )

                    if line_text:
                        paragraph_lines.append(
                            line_text
                        )

                if not paragraph_lines:
                    continue

                # =========================================
                # FONT ANALYSIS
                # =========================================

                first_span = lines[0]["spans"][0]

                font_size = round(
                    first_span.get(
                        "size",
                        12
                    )
                )

                flags = first_span.get(
                    "flags",
                    0
                )

                is_bold = bool(
                    flags & 16
                )

                # =========================================
                # ALIGNMENT
                # =========================================

                center_ratio = (
                    (bx0 + bx1) / 2
                ) / page_width

                is_center = (
                    0.40 <= center_ratio <= 0.60
                )

                # =========================================
                # HEADING DETECTION
                # =========================================

                is_heading = (

                    font_size >= (
                        common_size + 2
                    )

                    or (

                        is_bold
                        and len(
                            " ".join(
                                paragraph_lines
                            )
                        ) < 120
                    )
                )

                # =========================================
                # LIST DETECTION
                # =========================================

                first_line = paragraph_lines[0]

                is_list = (

                    first_line.startswith(
                        (
                            "•",
                            "-",
                            "*"
                        )
                    )

                    or re.match(
                        r'^\d+[\.\)]',
                        first_line
                    )
                )

                # =========================================
                # SPACING DETECTION
                # =========================================

                add_space_before = False

                if previous_y is not None:

                    y_gap = by0 - previous_y

                    if y_gap > font_size * 1.5:

                        add_space_before = True

                # =========================================
                # CREATE PARAGRAPH
                # =========================================

                if (
                    is_heading
                    and is_center
                ):

                    p = doc.add_paragraph(
                        style='Heading 1'
                    )

                    p.alignment = (
                        WD_ALIGN_PARAGRAPH.CENTER
                    )

                elif is_heading:

                    p = doc.add_paragraph(
                        style='Heading 2'
                    )

                    p.alignment = (
                        WD_ALIGN_PARAGRAPH.LEFT
                    )

                elif is_list:

                    p = doc.add_paragraph(
                        style='List Paragraph'
                    )

                else:

                    p = doc.add_paragraph()

                    # justify
                    if (
                        block_width >
                        page_width * 0.70
                    ):

                        p.alignment = (
                            WD_ALIGN_PARAGRAPH.JUSTIFY
                        )

                    else:

                        p.alignment = (
                            WD_ALIGN_PARAGRAPH.LEFT
                        )

                # =========================================
                # TYPOGRAPHY
                # =========================================

                p.paragraph_format.line_spacing = 1.15

                p.paragraph_format.space_after = Pt(6)

                if add_space_before:

                    p.paragraph_format.space_before = Pt(10)

                # =========================================
                # ADD TEXT
                # =========================================

                for line_index, line in enumerate(lines):

                    spans = line.get(
                        "spans",
                        []
                    )

                    if (
                        line_index > 0
                        and p.runs
                    ):

                        last_text = (
                            p.runs[-1].text
                        )

                        if (
                            last_text
                            and not last_text.endswith("-")
                        ):

                            p.add_run(" ")

                    for span in spans:

                        text = clean_text(
                            span.get(
                                "text",
                                ""
                            )
                        )

                        if not text:
                            continue

                        # hyphen fix
                        if (
                            p.runs
                            and p.runs[-1]
                            .text.endswith("-")
                        ):

                            p.runs[-1].text = (
                                p.runs[-1]
                                .text[:-1]
                            )

                            run = p.add_run(
                                text
                            )

                        else:

                            run = p.add_run(
                                text
                            )

                        # =====================================
                        # FONT
                        # =====================================

                        run.font.size = Pt(
                            round(
                                span.get(
                                    "size",
                                    12
                                )
                            )
                        )

                        flags = span.get(
                            "flags",
                            0
                        )

                        run.bold = bool(
                            flags & 16
                        )

                        run.italic = bool(
                            flags & 2
                        )

                        font_name = span.get(
                            "font",
                            "Times New Roman"
                        )

                        if "+" in font_name:

                            font_name = (
                                font_name
                                .split("+")[-1]
                            )

                        font_name = (
                            font_name.replace(
                                "-Identity-H",
                                ""
                            )
                        )

                        run.font.name = (
                            font_name
                        )

                        # =====================================
                        # COLOR
                        # =====================================

                        r, g, b = get_rgb(
                            span.get(
                                "color",
                                0
                            )
                        )

                        if not (
                            r < 25
                            and g < 25
                            and b < 25
                        ):

                            run.font.color.rgb = (
                                RGBColor(
                                    r,
                                    g,
                                    b
                                )
                            )

                    if line_index < len(lines) - 1:

                        p.add_run("\n")

                previous_y = by1

        pdf.close()

        # =================================================
        # SAVE
        # =================================================

        doc.save(str(output_path))

        print("=" * 60)
        print("SUCCESS ADVANCED RECONSTRUCTION")
        print("=" * 60)

    except Exception as e:

        print("=" * 60)
        print("SEMUA METHOD GAGAL")
        print("=" * 60)

        print(e)

# ─── KONVERSI: PPT → PDF ────────────────────────────────────────────────

def convert_ppt_to_pdf(input_file, output_path):
    output_dir = Path(output_path).parent
    result = run_libreoffice(input_file, 'pdf', output_dir)
    if str(result) != str(output_path):
        shutil.move(str(result), str(output_path))

# ─── KONVERSI: PDF → PPT ────────────────────────────────────────────────

def convert_pdf_to_ppt(input_file, output_path):
    import fitz
    import io
    import re
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    # =====================================================
    # SETUP PRESENTASI
    # =====================================================
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    SLIDE_W = prs.slide_width
    SLIDE_H = prs.slide_height
    blank_layout = prs.slide_layouts[6]

    pdf = fitz.open(str(input_file))

    def sx(x, pdf_w): return Emu(int((x / pdf_w) * SLIDE_W))
    def sy(y, pdf_h): return Emu(int((y / pdf_h) * SLIDE_H))
    def int_to_rgb(color_int):
        return RGBColor((color_int >> 16) & 255, (color_int >> 8) & 255, color_int & 255)

    # =====================================================
    # LOOP HALAMAN
    # =====================================================
    for page_index in range(len(pdf)):
        print(f"Convert halaman {page_index + 1}...")
        page = pdf[page_index]
        pdf_w = page.rect.width
        pdf_h = page.rect.height
        slide = prs.slides.add_slide(blank_layout)

        font_scale = 540.0 / pdf_h if pdf_h > 0 else 1.0

        text_dict = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE | fitz.TEXT_PRESERVE_LIGATURES)
        blocks = text_dict.get("blocks", [])

        # =================================================
        # HAPUS TEKS DARI PDF UNTUK BACKGROUND
        # =================================================
        for xref in page.get_contents():
            stream = pdf.xref_stream(xref)
            if stream:
                new_stream = re.sub(br'\bBT\b[\s\S]*?\bET\b', b'', stream)
                pdf.update_stream(xref, new_stream)
                
        for xobj in page.get_xobjects():
            xref = xobj[0]
            stream = pdf.xref_stream(xref)
            if stream:
                new_stream = re.sub(br'\bBT\b[\s\S]*?\bET\b', b'', stream)
                pdf.update_stream(xref, new_stream)

        # =================================================
        # RENDER BACKGROUND BERSIH (RAM OPTIMIZED)
        # =================================================
        # PERUBAHAN: Matriks 2x2. Hemat RAM, tapi tetap HD.
        matrix = fitz.Matrix(2, 2) 
        pix = page.get_pixmap(matrix=matrix, alpha=False)
        bg_stream = io.BytesIO(pix.tobytes("png"))
        slide.shapes.add_picture(bg_stream, 0, 0, width=SLIDE_W, height=SLIDE_H)

        # =================================================
        # TEKS EDITABLE (SPASI NORMAL & JUSTIFY)
        # =================================================
        for block in blocks:
            if block.get("type") != 0: continue
            
            bbox = block.get("bbox")
            if not bbox: continue

            lines = block.get("lines", [])
            if not lines: continue

            x0, y0, x1, y1 = bbox
            left = sx(x0, pdf_w)
            top = sy(y0, pdf_h)
            width = sx(x1 - x0, pdf_w) + Pt(12)
            height = sy(y1 - y0, pdf_h) + Pt(12)

            if width < Emu(10000) or height < Emu(10000): continue

            # Deteksi Alignment Rata Kiri-Kanan (Justify)
            x0_list = [l.get("bbox", [0,0,0,0])[0] for l in lines]
            x1_list = [l.get("bbox", [0,0,0,0])[2] for l in lines]
            
            if len(lines) == 1:
                cx = (x0_list[0] + x1_list[0]) / 2
                if abs(cx - pdf_w/2) < (pdf_w * 0.05): align = PP_ALIGN.CENTER
                else: align = PP_ALIGN.LEFT
            else:
                diff_x0 = max(x0_list) - min(x0_list)
                diff_x1 = max(x1_list) - min(x1_list)
                
                if diff_x0 < 25 and diff_x1 < 25: align = PP_ALIGN.JUSTIFY
                elif diff_x0 < 25: align = PP_ALIGN.LEFT
                elif diff_x1 < 25: align = PP_ALIGN.RIGHT
                else: align = PP_ALIGN.CENTER

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.clear()
            tf.word_wrap = True
            tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

            p = tf.paragraphs[0]
            p.alignment = align

            for i, line in enumerate(lines):
                if i > 0:
                    r_space = p.add_run()
                    r_space.text = " "
                    r_space.font.size = Pt(12) 

                for span in line.get("spans", []):
                    text = span.get("text", "")
                    
                    if not text: continue 

                    run = p.add_run()
                    run.text = text

                    sz = span.get("size", 12) * font_scale
                    run.font.size = Pt(max(6, min(sz, 96)))

                    font_name = span.get("font", "")
                    if "+" in font_name: font_name = font_name.split("+")[-1]
                    font_name = font_name.replace("-Identity-H", "").split(",")[0]
                    if font_name: run.font.name = font_name

                    flags = span.get("flags", 0)
                    run.font.bold = bool(flags & 16)
                    run.font.italic = bool(flags & 2)

                    run.font.color.rgb = int_to_rgb(span.get("color", 0))

    pdf.close()
    prs.save(str(output_path))
    print("Selesai convert!")

# ─── KONVERSI: Word → PPT ────────────────────────────────────────────────

def convert_word_to_ppt(input_file, output_path):
    from docx import Document
    from docx.document import Document as _Document
    from docx.oxml.text.paragraph import CT_P
    from docx.oxml.table import CT_Tbl
    from docx.table import _Cell, Table
    from docx.text.paragraph import Paragraph
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    import re

    # ─── FUNGSI HELPER: Membaca Paragraf & Tabel Sesuai Urutan ───
    def iter_block_items(parent):
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            return
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    doc = Document(str(input_file))
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    H1_STYLES = {'heading 1', 'heading1', 'title', 'judul'}
    H2_STYLES = {'heading 2', 'heading2', 'subtitle', 'subjudul'}

    slides_data = []
    cur_slide = {'title': None, 'subtitle': None, 'content': []}
    last_h1_title = "Data Presentasi" 

    def save_current_slide():
        nonlocal cur_slide
        if cur_slide['title'] or cur_slide['content']:
            slides_data.append(cur_slide)
        cur_slide = {'title': None, 'subtitle': None, 'content': []}

    # ─── 1. PARSING CERDAS (Ekstrak Teks & TABEL) ───
    for block in iter_block_items(doc):
        # PENANGANAN TEKS PARAGRAF
        if isinstance(block, Paragraph):
            text = block.text.strip()
            if not text: continue
            
            style_name = block.style.name.lower().replace(' ', '')
            
            runs_data = []
            for r in block.runs:
                if r.text.strip():
                    runs_data.append({
                        'text': r.text,
                        'bold': r.bold if r.bold is not None else False,
                        'italic': r.italic if r.italic is not None else False
                    })

            if any(h in style_name for h in H1_STYLES):
                save_current_slide()
                cur_slide['title'] = text
                last_h1_title = text
            elif any(h in style_name for h in H2_STYLES):
                if not cur_slide['title']:
                    save_current_slide()
                    cur_slide['title'] = text
                else:
                    cur_slide['subtitle'] = text
            else:
                is_bullet = bool(re.match(r'^[\u2022\u2023•▪▸►\-\*]\s', text))
                is_number = bool(re.match(r'^\d+[.)]\s', text))
                cur_slide['content'].append({
                    'type': 'bullet' if (is_bullet or is_number) else 'body',
                    'raw_text': text,
                    'runs': runs_data
                })

        # PENANGANAN TABEL
        elif isinstance(block, Table):
            if cur_slide['content']:
                save_current_slide() # Pisahkan teks sebelumnya ke slide sendiri
            
            table_data = [[cell.text.strip() for cell in row.cells] for row in block.rows]
            
            if table_data:
                header = table_data[0]
                body = table_data[1:]
                
                # Fitur Pintar: Memecah tabel panjang jadi beberapa slide (7 baris/slide)
                chunk_size = 7 
                
                if not body:
                    cur_slide['title'] = f"{last_h1_title} (Tabel)"
                    cur_slide['content'] = [{'type': 'table', 'data': [header]}]
                    save_current_slide()
                else:
                    for i in range(0, len(body), chunk_size):
                        chunk = body[i:i+chunk_size]
                        cur_slide['title'] = f"{last_h1_title} (Lanjutan)" if i > 0 else f"{last_h1_title} (Tabel)"
                        cur_slide['content'] = [{'type': 'table', 'data': [header] + chunk}]
                        save_current_slide()

    save_current_slide()

    if not slides_data:
        import os
        slides_data = [{'title': os.path.basename(input_file), 'subtitle': 'Dokumen Kosong / Tidak Terbaca', 'content': []}]

    # ─── 2. AUTO-DESIGN ENGINE (Menggambar Slide PPTX) ───
    COLOR_BG     = RGBColor(15, 27, 56)
    COLOR_ACCENT = RGBColor(58, 134, 255)
    COLOR_TITLE  = RGBColor(255, 255, 255)
    COLOR_BODY   = RGBColor(208, 228, 255)
    COLOR_MUTED  = RGBColor(112, 144, 192)

    def set_background(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def draw_rect(slide, left, top, width, height, color):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()

    for index, sd in enumerate(slides_data):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_background(slide, COLOR_BG)
        draw_rect(slide, 0, 0, Inches(0.15), prs.slide_height, COLOR_ACCENT)

        has_content = bool(sd['content'])

        if has_content:
            # === RENDER JUDUL ===
            title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(1.2))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = sd['title'] or "Slide Data"
            r.font.bold = True
            r.font.size = Pt(36)
            r.font.color.rgb = COLOR_TITLE

            draw_rect(slide, Inches(0.4), Inches(1.4), Inches(12.0), Pt(2), COLOR_ACCENT)

            # Cek apakah ini Slide Tabel atau Slide Teks
            is_table_slide = len(sd['content']) == 1 and sd['content'][0]['type'] == 'table'

            if is_table_slide:
                # === RENDER TABEL ASLI ===
                table_data = sd['content'][0]['data']
                rows = len(table_data)
                cols = max(len(row) for row in table_data) if rows > 0 else 0
                
                if rows > 0 and cols > 0:
                    tbl_shape = slide.shapes.add_table(rows, cols, Inches(0.4), Inches(1.7), Inches(12.0), Inches(0.6 * rows))
                    table = tbl_shape.table
                    
                    for r_idx, row in enumerate(table_data):
                        for c_idx, cell_text in enumerate(row):
                            if c_idx < cols:
                                cell = table.cell(r_idx, c_idx)
                                cell.text = cell_text
                                
                                # Desain Teks Dalam Tabel
                                for paragraph in cell.text_frame.paragraphs:
                                    paragraph.font.size = Pt(14)
                                    paragraph.font.color.rgb = COLOR_TITLE if r_idx == 0 else RGBColor(0,0,0)
                                    paragraph.font.bold = (r_idx == 0)
                                    paragraph.alignment = PP_ALIGN.CENTER if r_idx == 0 else PP_ALIGN.LEFT
                                    
                                # Desain Warna Latar Sel Tabel
                                cell.fill.solid()
                                if r_idx == 0:
                                    cell.fill.fore_color.rgb = COLOR_ACCENT # Header warna biru terang
                                elif r_idx % 2 == 0:
                                    cell.fill.fore_color.rgb = RGBColor(220, 235, 255) # Belang-belang rapi
                                else:
                                    cell.fill.fore_color.rgb = RGBColor(240, 248, 255)
            else:
                # === RENDER TEKS BIASA ===
                content_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(12.5), Inches(5.4))
                ctf = content_box.text_frame
                ctf.word_wrap = True

                for j, item in enumerate(sd['content'][:10]): 
                    cp = ctf.paragraphs[0] if j == 0 else ctf.add_paragraph()
                    cp.space_before = Pt(8)
                    
                    if item['type'] == 'bullet':
                        dot = cp.add_run()
                        dot.text = "▪  "
                        dot.font.color.rgb = COLOR_ACCENT
                        dot.font.bold = True
                    
                    for run_data in item['runs']:
                        clean_text = re.sub(r'^[\u2022\u2023•▪▸►\-\*\d+.)\s]+', '', run_data['text']) if item['type'] == 'bullet' else run_data['text']
                        if not clean_text: continue
                        
                        r_body = cp.add_run()
                        r_body.text = clean_text
                        r_body.font.bold = run_data['bold']
                        r_body.font.italic = run_data['italic']
                        r_body.font.size = Pt(22)
                        r_body.font.color.rgb = COLOR_BODY

        else:
            # === RENDER COVER ===
            tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.5), Inches(11.33), Inches(3.0))
            tf = tb.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            r = p.add_run()
            r.text = sd['title'] or "Presentasi"
            r.font.bold = True
            r.font.size = Pt(54)
            r.font.color.rgb = COLOR_TITLE

            if sd['subtitle']:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                p2.space_before = Pt(10)
                r2 = p2.add_run()
                r2.text = sd['subtitle']
                r2.font.size = Pt(28)
                r2.font.color.rgb = COLOR_MUTED

        # Nomor Slide
        num_box = slide.shapes.add_textbox(Inches(12.0), Inches(6.9), Inches(1.0), Inches(0.5))
        np = num_box.text_frame.paragraphs[0]
        np.alignment = PP_ALIGN.RIGHT
        n_run = np.add_run()
        n_run.text = f"{index + 1}"
        n_run.font.size = Pt(14)
        n_run.font.color.rgb = COLOR_MUTED

    prs.save(str(output_path))

# ─── KONVERSI: Excel → PDF ───────────────────────────────────────────────

def convert_excel_to_pdf(input_file, output_path):

    import os
    import tempfile
    import shutil
    import subprocess

    from openpyxl import load_workbook

    # =====================================================
    # TEMP FILE
    # =====================================================

    temp_dir = tempfile.mkdtemp()

    fixed_xlsx = os.path.join(
        temp_dir,
        "fixed.xlsx"
    )

    # =====================================================
    # LOAD WORKBOOK
    # =====================================================

    wb = load_workbook(input_file)

    # =====================================================
    # FIX SEMUA SHEET
    # =====================================================

    for ws in wb.worksheets:

        # =============================================
        # PAGE SETUP
        # =============================================

        ws.page_setup.orientation = 'landscape'

        ws.page_setup.paperSize = ws.PAPERSIZE_A4

        ws.page_setup.fitToWidth = 1
        ws.page_setup.fitToHeight = 0

        ws.sheet_properties.pageSetUpPr.fitToPage = True

        # =============================================
        # MARGIN
        # =============================================

        ws.page_margins.left = 0.2
        ws.page_margins.right = 0.2

        ws.page_margins.top = 0.3
        ws.page_margins.bottom = 0.3

        # =============================================
        # CENTER PAGE
        # =============================================

        ws.print_options.horizontalCentered = True

        # =============================================
        # AUTO WIDTH
        # =============================================

        for column_cells in ws.columns:

            length = 0

            column = column_cells[0].column_letter

            for cell in column_cells:

                try:

                    if cell.value:

                        length = max(
                            length,
                            len(str(cell.value))
                        )

                except:
                    pass

            adjusted_width = min(
                length + 3,
                50
            )

            ws.column_dimensions[
                column
            ].width = adjusted_width

        # =============================================
        # AUTO HEIGHT
        # =============================================

        for row in ws.iter_rows():

            max_lines = 1

            row_number = row[0].row

            for cell in row:

                if cell.value:

                    lines = str(
                        cell.value
                    ).count("\n") + 1

                    max_lines = max(
                        max_lines,
                        lines
                    )

            ws.row_dimensions[
                row_number
            ].height = max(
                20,
                max_lines * 15
            )

        # =============================================
        # PRINT AREA
        # =============================================

        max_row = ws.max_row
        max_col = ws.max_column

        from openpyxl.utils import (
            get_column_letter
        )

        last_col = get_column_letter(
            max_col
        )

        ws.print_area = (
            f"A1:{last_col}{max_row}"
        )

    # =====================================================
    # SAVE FIXED FILE
    # =====================================================

    wb.save(fixed_xlsx)

    # =====================================================
    # CONVERT VIA LIBREOFFICE
    # =====================================================

    output_dir = os.path.dirname(
        os.path.abspath(output_path)
    )

    cmd = [

    r"C:\Program Files\LibreOffice\program\soffice.exe",

    "--headless",

    "--convert-to",
    "pdf",

    fixed_xlsx,

    "--outdir",
    output_dir
    ]

    subprocess.run(
        cmd,
        check=True
    )

    # =====================================================
    # MOVE RESULT
    # =====================================================

    generated_pdf = os.path.join(

        output_dir,

        os.path.splitext(
            os.path.basename(
                fixed_xlsx
            )
        )[0] + ".pdf"
    )

    if generated_pdf != output_path:

        shutil.move(
            generated_pdf,
            output_path
        )

    print("SUCCESS EXCEL TO PDF")

# ─── KONVERSI: Gambar → PPT ──────────────────────────────────────────────

def convert_images_to_ppt(input_files, output_path):
    from pptx import Presentation
    from pptx.util import Inches, Emu
    from PIL import Image
    import io

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for img_path in input_files:
        img = Image.open(str(img_path))
        if img.mode in ('RGBA', 'LA', 'P'):
            img = img.convert('RGB')

        img_w, img_h = img.size
        slide_w = prs.slide_width
        slide_h = prs.slide_height

        # Scale gambar agar muat & proporsional, posisikan di tengah
        scale = min(slide_w / img_w, slide_h / img_h)
        new_w = int(img_w * scale)
        new_h = int(img_h * scale)
        left  = (slide_w - new_w) // 2
        top   = (slide_h - new_h) // 2

        slide = prs.slides.add_slide(blank)
        buf = io.BytesIO()
        img.save(buf, 'JPEG', quality=90)
        buf.seek(0)
        slide.shapes.add_picture(buf, left, top, new_w, new_h)

    prs.save(str(output_path))
    
# compres file
def compress_universal(input_files, output_path):
    import os
    import fitz
    import zipfile
    import shutil
    import tempfile
    from PIL import Image

    # PERBAIKAN: Ambil elemen pertama dari list dan ubah menjadi string bersih
    # Ini akan menghilangkan error [WindowsPath(...)]
    input_file = str(input_files[0])
    output_path = str(output_path)

    ext = os.path.splitext(input_file)[1].lower()

    try:
        # =====================================================
        # JALUR 1: KOMPRESI GAMBAR NATIVE
        # =====================================================
        if ext in ['.jpg', '.jpeg', '.png', '.webp', '.bmp']:
            img = Image.open(input_file)
            if img.mode in ('RGBA', 'LA', 'P') or 'transparency' in img.info:
                img.save(output_path, optimize=True)
            else:
                if img.mode != 'RGB':
                    img = img.convert('RGB')
                img.save(output_path, "JPEG", quality=50, optimize=True)

        # =====================================================
        # JALUR 2: KOMPRESI PDF 
        # =====================================================
        elif ext == '.pdf':
            doc = fitz.open(input_file)
            doc.save(output_path, garbage=4, deflate=True, clean=True)
            doc.close()

        # =====================================================
        # JALUR 3: DEEP COMPRESSION OFFICE (Word, Excel, PPT)
        # =====================================================
        elif ext in ['.docx', '.pptx', '.xlsx']:
            with tempfile.TemporaryDirectory() as tmpdir:
                with zipfile.ZipFile(input_file, 'r') as zip_ref:
                    zip_ref.extractall(tmpdir)

                for root, dirs, files in os.walk(tmpdir):
                    if 'media' in root: 
                        for f in files:
                            f_ext = os.path.splitext(f)[1].lower()
                            if f_ext in ['.jpg', '.jpeg', '.png']:
                                img_path = os.path.join(root, f)
                                try:
                                    img = Image.open(img_path)
                                    if img.mode != 'RGB':
                                        img = img.convert('RGB')
                                    img.save(img_path, "JPEG", quality=50, optimize=True)
                                except:
                                    pass 

                with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED, compresslevel=9) as zip_out:
                    for root, dirs, files in os.walk(tmpdir):
                        for file in files:
                            file_path = os.path.join(root, file)
                            arcname = os.path.relpath(file_path, tmpdir)
                            zip_out.write(file_path, arcname)

        # =====================================================
        # JALUR 4: FALLBACK (Semua file lainnya, divakum jadi ZIP)
        # =====================================================
        else:
            zip_path = output_path + ".zip"
            with zipfile.ZipFile(zip_path, 'w', zipfile.ZIP_DEFLATED, compresslevel=9) as zipf:
                zipf.write(input_file, arcname=os.path.basename(input_file))
            shutil.move(zip_path, output_path)
            
    except Exception as e:
        raise RuntimeError(f"Gagal mengompres file: {str(e)}")

# ─── CONFIG TOOL ─────────────────────────────────────────────────────────

TOOLS = {
    'image-to-pdf': {
        'label': 'Gambar ke PDF',
        'input_exts': {'.jpg','.jpeg','.png','.gif','.bmp','.tiff','.tif','.webp'},
        'output_ext': '.pdf',
        'multi_input': True,
        'fn': lambda ins, out: convert_images_to_pdf(ins, out),
    },
    'image-to-word': {
        'label': 'Gambar ke Word',
        'input_exts': {'.jpg','.jpeg','.png','.gif','.bmp','.tiff','.tif','.webp'},
        'output_ext': '.docx',
        'multi_input': True,
        'fn': lambda ins, out: convert_images_to_word(ins, out),
    },
    'word-to-pdf': {
        'label': 'Word ke PDF',
        'input_exts': {'.doc','.docx'},
        'output_ext': '.pdf',
        'multi_input': False,
        'fn': lambda ins, out: convert_word_to_pdf(ins[0], out),
    },
    'pdf-to-word': {
        'label': 'PDF ke Word',
        'input_exts': {'.pdf'},
        'output_ext': '.docx',
        'multi_input': False,
        'fn': lambda ins, out: convert_pdf_to_word(ins[0], out),
    },
    'ppt-to-pdf': {
        'label': 'PPT ke PDF',
        'input_exts': {'.ppt','.pptx'},
        'output_ext': '.pdf',
        'multi_input': False,
        'fn': lambda ins, out: convert_ppt_to_pdf(ins[0], out),
    },
    'pdf-to-ppt': {
        'label': 'PDF ke PPT',
        'input_exts': {'.pdf'},
        'output_ext': '.pptx',
        'multi_input': False,
        'fn': lambda ins, out: convert_pdf_to_ppt(ins[0], out),
    },
    'word-to-ppt': {
        'label': 'Word ke PPT',
        'input_exts': {'.doc','.docx'},
        'output_ext': '.pptx',
        'multi_input': False,
        'fn': lambda ins, out: convert_word_to_ppt(ins[0], out),
    },
    'excel-to-pdf': {
        'label': 'Excel ke PDF',
        'input_exts': {'.xls','.xlsx','.ods'},
        'output_ext': '.pdf',
        'multi_input': False,
        'fn': lambda ins, out: convert_excel_to_pdf(ins[0], out),
    },
    'image-to-ppt': {
        'label': 'Gambar ke PPT',
        'input_exts': {'.jpg','.jpeg','.png','.gif','.bmp','.tiff','.tif','.webp'},
        'output_ext': '.pptx',
        'multi_input': True,
        'fn': lambda ins, out: convert_images_to_ppt(ins, out),
    },
    'compress-file': {
        'fn': compress_universal,
        'input_exts': '*', # Tanda bintang = Bebas semua jenis file
        'output_ext': 'dynamic', # Ekstensi outputnya bunglon (berubah-ubah)
        'multi_input': False
    },
}

# ─── ROUTES ──────────────────────────────────────────────────────────────

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/api/convert', methods=['POST'])
def convert_api():
    conv_type = request.form.get('type', '').strip()
    if conv_type not in TOOLS:
        return jsonify({'error': f'Tipe konversi tidak valid: {conv_type}'}), 400

    tool   = TOOLS[conv_type]
    files  = request.files.getlist('files')
    if not files or all(f.filename == '' for f in files):
        return jsonify({'error': 'Tidak ada file yang dikirim'}), 400

    job_id  = str(uuid.uuid4())
    job_dir = TEMP_DIR / job_id
    job_dir.mkdir(parents=True, exist_ok=True)

    # ── Simpan file upload ──
    saved = []
    for f in files:
        fname = secure_filename(f.filename) or f'file_{uuid.uuid4().hex[:8]}'
        ext = Path(fname).suffix.lower()
        
        # PERUBAHAN: Loloskan jika input_exts adalah '*'
        if tool['input_exts'] != '*' and ext not in tool['input_exts']:
            shutil.rmtree(str(job_dir), ignore_errors=True)
            return jsonify({'error': f'Format tidak didukung: {ext}.'}), 400
            
        dest = job_dir / fname
        f.save(str(dest))
        saved.append(dest)

    # ── Jalankan konversi ──
    output_files = []
    try:
        if tool['multi_input']:
            out_name = 'converted' + tool['output_ext']
            out_path = job_dir / out_name
            tool['fn'](saved, str(out_path))
            output_files.append({'name': out_name, 'path': str(out_path)})
        else:
            for f_path in saved:
                # PERUBAHAN: Logika Ekstensi Dinamis untuk Kompresor
                if tool['output_ext'] == 'dynamic':
                    asli_ext = f_path.suffix.lower()
                    # Jika file didukung kompresi native, pertahankan ekstensinya
                    if asli_ext in ['.jpg', '.jpeg', '.png', '.webp', '.bmp', '.pdf', '.docx', '.pptx', '.xlsx']:
                        out_ext = asli_ext
                    else:
                        out_ext = '.zip' # Jika file asing, jadikan ZIP
                    
                    out_name = f_path.stem + '_compressed' + out_ext
                else:
                    out_name = f_path.stem + tool['output_ext']
                
                out_path = job_dir / out_name
                tool['fn']([f_path], str(out_path))
                output_files.append({'name': out_name, 'path': str(out_path)})

    except Exception as e:
        shutil.rmtree(str(job_dir), ignore_errors=True)
        return jsonify({'error': str(e)}), 500

    # Verifikasi output ada
    for of in output_files:
        if not Path(of['path']).exists():
            shutil.rmtree(str(job_dir), ignore_errors=True)
            return jsonify({'error': f'File output tidak ditemukan: {of["name"]}'}), 500

    jobs[job_id] = {
        'dir': str(job_dir),
        'files': output_files,
        'created': time.time()
    }

    return jsonify({
        'success': True,
        'job_id': job_id,
        'files': [{'name': f['name']} for f in output_files]
    })

@app.route('/api/download/<job_id>/<path:filename>')
def download_file(job_id, filename):
    if job_id not in jobs:
        return jsonify({'error': 'Job tidak ditemukan atau sudah kadaluarsa'}), 404
    job = jobs[job_id]
    file_path = Path(job['dir']) / filename
    if not file_path.exists():
        return jsonify({'error': 'File tidak ditemukan'}), 404
    return send_file(str(file_path), as_attachment=True, download_name=filename)

@app.route('/api/download-zip/<job_id>')
def download_zip(job_id):
    if job_id not in jobs:
        return jsonify({'error': 'Job tidak ditemukan'}), 404
    job    = jobs[job_id]
    job_dir = Path(job['dir'])
    zip_path = job_dir / 'results.zip'
    with zipfile.ZipFile(str(zip_path), 'w', zipfile.ZIP_DEFLATED) as zf:
        for f in job['files']:
            fp = Path(f['path'])
            if fp.exists():
                zf.write(str(fp), f['name'])
    return send_file(str(zip_path), as_attachment=True, download_name='converted_files.zip')

@app.route('/api/status')
def status():
    lo = get_libreoffice()
    return jsonify({
        'libreoffice': lo is not None,
        'libreoffice_path': lo,
        'poppler': get_poppler_path() is not None or shutil.which('pdftoppm') is not None,
        'platform': platform.system(),
        'jobs_active': len(jobs)
    })

# ─── CLEANUP ─────────────────────────────────────────────────────────────

def cleanup_loop():
    """Hapus job yang sudah lebih dari 2 jam"""
    while True:
        time.sleep(1800)
        now = time.time()
        expired = [jid for jid, j in list(jobs.items()) if now - j['created'] > 7200]
        for jid in expired:
            job = jobs.pop(jid, None)
            if job:
                shutil.rmtree(job['dir'], ignore_errors=True)

threading.Thread(target=cleanup_loop, daemon=True).start()

# ─── MAIN ─────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    print("=" * 55)
    print("  🚀  LocalConvert - Konversi File Lokal")
    print("=" * 55)
    print(f"  Buka browser: http://localhost:5000")
    print(f"  LibreOffice : {get_libreoffice() or '❌ Tidak ditemukan'}")
    print(f"  Temp dir    : {TEMP_DIR}")
    print("=" * 55)
    app.run(debug=True, host='0.0.0.0', port=5000, threaded=True)
